from langchain.tools import tool
#from langchain_ollama import ChatOllama
from langchain_groq import ChatGroq
from langgraph.graph import StateGraph, MessagesState, START, END
from langchain.messages import AnyMessage, HumanMessage, AIMessage
from typing_extensions import TypedDict, Annotated
import operator
from langchain.messages import SystemMessage
from langchain.messages import ToolMessage
from typing import Literal
from langchain_core.language_models import BaseChatModel
import json
import re
import base64
import os
import shutil
from RAG import RAGEngine

from pptx import Presentation
from pptx.util import Inches


class SlideGenerationGraphState(TypedDict):
    messages:Annotated[list[str], operator.add]
    llm_calls:int

    #user_question:str
    retrieved_text:str
    retrieved_images:list[str] # <-- ?????

    Text_Summarizer_output:str
    Image_Captioner_output:str
    Code_Generator_output:str # without review

    # --- ADD THIS FIELD ---
    json_presentation_data: str # Stores the structured JSON for the PPTX renderer
    # ----------------------
    
    Code_Reviewer_output:str
    Page_Reviewer_output:str
    Code_Generator_output_Reviewed:str
    
class SlideGenerationModule():
    def __init__(self, retriever:RAGEngine, model:BaseChatModel=ChatGroq(model="meta-llama/llama-4-scout-17b-16e-instruct")):
        self.__retriever = retriever
        self.__model = model
        self.__workflow_builder = StateGraph(SlideGenerationGraphState)
        
        def Text_Summarizer(state:dict)   :
            agent_answer:AIMessage = self.__model.invoke(
                [
                    SystemMessage(
                            content=TS_SYSTEM_PROMPT
                        ),
                    HumanMessage(
                            content=TS_USER_PROMPT.replace(
                                "<Document text>",
                                state["retrieved_text"]
                            )
                        )
                ]
            )
            return {
                "messages": [agent_answer.content],
                "llm_calls": state.get('llm_calls', 0) + 1,
                "Text_Summarizer_output": agent_answer.content
            }
            
        def Image_Captioner(state: dict):
            # 1. Get Image Filenames to show the LLM
            retrieved_imgs = state.get("retrieved_images", [])
            
            # Create a string list of filenames: "Image 1: chart.png", etc.
            filename_list = []
            for i, path in enumerate(retrieved_imgs):
                filename = os.path.basename(path)
                filename_list.append(f"Image {i+1} Filename: {filename}")
            
            filenames_str = "\n".join(filename_list)

            # 2. Update the prompt to include these filenames
            text_prompt = IC_USER_PROMPT.replace("<Document Text>", state["retrieved_text"])
            text_prompt += f"\n\nREFERENCED IMAGE FILENAMES:\n{filenames_str}"
            
            # 3. Build the message (using the fix from before)
            message_content = [{"type": "text", "text": text_prompt}]
            
            for img_path in retrieved_imgs:
                # Use the helper to get base64 (from previous fix)
                base64_image = self._encode_image(img_path)
                if base64_image:
                    message_content.append({
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}
                    })

            # 4. Invoke LLM
            agent_answer: AIMessage = self.__model.invoke([
                SystemMessage(content=IC_SYSTEM_PROMPT),
                HumanMessage(content=message_content)
            ])
            
            return {
                "messages": [agent_answer.content],
                "llm_calls": state.get('llm_calls', 0) + 1,
                "Image_Captioner_output": agent_answer.content
            }
            
        def Code_Generator(state:dict)   :
            agent_answer:AIMessage = self.__model.invoke(
                [
                    SystemMessage(
                            content=CG_SYSTEM_PROMPT
                        ),
                    HumanMessage(
                            content=CG_USER_PROMPT.replace(
                                "<SlidevGrammar>",
                                SLIDEV_GRAMMAR
                            ).replace(
                                "<TextSummary.md>",
                                state["Text_Summarizer_output"]
                            ).replace(
                                "<ImageCaption.md>",
                                state["Image_Captioner_output"]
                            )
                    )
                ]
            )
            return {
                "messages": [agent_answer.content],
                "llm_calls": state.get('llm_calls', 0) + 1,
                "Code_Generator_output": agent_answer.content
            }
            
        def Code_Reviewer(state:dict)   :
            agent_answer:AIMessage = self.__model.invoke(
                [
                    SystemMessage(
                            content=CR_SYSTEM_PROMPT
                        ),
                    HumanMessage(
                            content=CR_USER_PROMPT.replace(
                                "<SlidevGrammar>",
                                SLIDEV_GRAMMAR
                            ).replace(
                                "<TextSummary.md>",
                                state["Text_Summarizer_output"]
                            ).replace(
                                "<ImageCaption.md>",
                                state["Image_Captioner_output"]
                            ).replace(
                                "<SlidevCode.md>",
                                state["Code_Generator_output"]
                            )
                    )
                ]
            )
            return {
                "messages": [agent_answer.content],
                "llm_calls": state.get('llm_calls', 0) + 1,
                "Code_Reviewer_output": agent_answer.content
            }
            
        def Page_Reviewer(state:dict)   :
            agent_answer:AIMessage = self.__model.invoke(
                [
                    SystemMessage(
                            content=PR_SYSTEM_PROMPT
                        ),
                    HumanMessage(
                            content=PR_USER_PROMPT.replace(
                                "<SlidevPages>",
                                SLIDEV_PAGES
                            ).replace(
                                "<ImageCaption.md>",
                                state["Image_Captioner_output"]
                            ).replace(
                                "<Document Images>",
                                "\n".join(state["retrieved_images"])
                            )
                    )
                ]
            )
            return {
                "messages": [agent_answer.content],
                "llm_calls": state.get('llm_calls', 0) + 1,
                "Page_Reviewer_output": agent_answer.content
            }
            
        def Code_Generator_Reviewed(state:dict)   :
            agent_answer:AIMessage = self.__model.invoke(
                [
                    SystemMessage(
                            content=CGR_SYSTEM_PROMPT
                        ),
                    HumanMessage(
                            content=CGR_USER_PROMPT.replace(
                                "<SlidevGrammar>",
                                SLIDEV_GRAMMAR
                            ).replace(
                                "<TextSummary.md>",
                                state["Text_Summarizer_output"]
                            ).replace(
                                "<ImageCaption.md>",
                                state["Image_Captioner_output"]
                            ).replace(
                                "<CodeReview.md>",
                                state["Code_Reviewer_output"]
                            )
                    )
                ]
            )
            return {
                "messages": [agent_answer.content],
                "llm_calls": state.get('llm_calls', 0) + 1,
                "Code_Generator_output_Reviewed": agent_answer.content
            }
            
        self.__workflow_builder.add_node("Text_Summarizer", Text_Summarizer)
        self.__workflow_builder.add_node("Image_Captioner", Image_Captioner)
        self.__workflow_builder.add_node("Code_Generator", Code_Generator)
        self.__workflow_builder.add_node("Code_Reviewer", Code_Reviewer)
        self.__workflow_builder.add_node("Page_Reviewer", Page_Reviewer)
        self.__workflow_builder.add_node("Code_Generator_Reviewed", Code_Generator_Reviewed)
        
        self.__workflow_builder.add_edge(START, "Text_Summarizer")
        self.__workflow_builder.add_edge("Text_Summarizer", "Image_Captioner")
        self.__workflow_builder.add_edge("Image_Captioner", "Code_Generator")
        self.__workflow_builder.add_edge("Code_Generator", "Code_Reviewer")
        self.__workflow_builder.add_edge("Code_Reviewer", "Page_Reviewer")
        self.__workflow_builder.add_edge("Page_Reviewer", "Code_Generator_Reviewed")
        self.__workflow_builder.add_edge("Code_Generator_Reviewed", END)
        
        self.__workflow = self.__workflow_builder.compile()
        
    # In SlideGeneration.py inside SlideGenerationModule class

    def generate_slides(self, topic: str):
        """
        Main entry point: Queries RAG, copies images to local folder, then runs generation.
        """
        print(f"Retrieving content for: {topic}...")
        rag_result = self.__retriever.query(topic, k_text=5, k_image=3)
        
        text_content = "\n\n".join(rag_result.get("text", []))
        original_images = rag_result.get("images", [])
        
        # --- CRITICAL FIX: Copy images to local 'images/' folder ---
        # This fixes the absolute path/backslash error in Slidev
        local_image_paths = []
        
        if original_images:
            os.makedirs("images", exist_ok=True) # Create folder
            print(f"Copying {len(original_images)} images to local 'images/' folder...")
            
            for src_path in original_images:
                if os.path.exists(src_path):
                    filename = os.path.basename(src_path)
                    # Sanitize filename (remove spaces, etc. if needed)
                    filename = filename.replace(" ", "_")
                    
                    dst_path = os.path.join("images", filename)
                    shutil.copy(src_path, dst_path)
                    
                    # Store as FORWARD SLASH path for Markdown
                    # e.g., "images/my_diagram.png"
                    local_image_paths.append(f"images/{filename}")
                else:
                    print(f"Warning: Source image not found: {src_path}")
        
        if not text_content:
            text_content = "No specific text found in documents."
            print("Warning: No text retrieved from RAG.")

        initial_state = {
            "messages": [HumanMessage(content=topic)],
            "llm_calls": 0,
            "retrieved_text": text_content,
            "retrieved_images": local_image_paths, # Pass the CLEAN local paths
            "Text_Summarizer_output": "",
            "Image_Captioner_output": "",
            "Code_Generator_output": "",
            "Code_Reviewer_output": "",
            "Page_Reviewer_output": "",
            "Code_Generator_output_Reviewed": ""
        }
        
        print("Starting Multi-Agent Workflow...")
        final_state = self.__workflow.invoke(initial_state)
        return final_state["Code_Generator_output_Reviewed"]
    
    def _encode_image(self, image_path):
        """Helper to check if valid file and convert to base64"""
        if not image_path or not isinstance(image_path, str):
            print(f"Warning: Invalid image path: {image_path}")
            return None
        
        if not os.path.exists(image_path):
            print(f"Warning: Image file not found: {image_path}")
            return None
            
        try:
            with open(image_path, "rb") as image_file:
                return base64.b64encode(image_file.read()).decode('utf-8')
        except Exception as e:
            print(f"Error encoding image {image_path}: {e}")
            return None
        
    def visualize_workflow(self):
        from IPython.display import Image, display
        display(Image(self.__workflow.get_graph().draw_mermaid_png()))
    

    def clean_json_string(self, s):
        print(f"--- [DEBUG] Attempting to extract JSON from content (Length: {len(s)})")
        
        # 1. Try to find content between the first [ and the last ]
        match = re.search(r'(\[.*\])', s, re.DOTALL)
        if match:
            s = match.group(1)
        
        # 2. Basic cleanup of markdown markers
        s = re.sub(r'```json\s*', '', s)
        s = re.sub(r'```', '', s)
        
        return s.strip()


    def save_as_pptx(self, raw_llm_output, output_path="GP_Presentation.pptx"):
        print(f"\n--- [RENDER] Starting PPTX Generation ---")
        prs = Presentation()
        
        cleaned_data = self.clean_json_string(raw_llm_output)
        
        try:
            slides_data = json.loads(cleaned_data)
            
            for i, data in enumerate(slides_data):
                # Layout 1 is 'Title and Content'
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                
                # 1. Set the Title
                slide.shapes.title.text = data.get("title", f"Slide {i+1}")
                
                # 2. Add Bullet Points (ensure we don't include paths here)
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear() # Remove default text
                
                for point in data.get("content", []):
                    p = tf.add_paragraph()
                    p.text = str(point)
                
                # 3. FIX: Actually Insert the Image (don't just print the path)
                img_path = data.get("image_path")
                if img_path and os.path.exists(img_path):
                    print(f"[RENDER] Inserting image: {img_path}")
                    # Position: Right side of the slide
                    # (left, top, width, height)
                    slide.shapes.add_picture(
                        img_path, 
                        left=Inches(6.5), 
                        top=Inches(1.5), 
                        width=Inches(3)
                    )
                else:
                    print(f"[WARNING] Image path invalid or missing: {img_path}")

            prs.save(output_path)
            print(f"--- [SUCCESS] Presentation saved to {output_path}")

        except Exception as e:
            print(f"--- [RENDER ERROR]: {e}")
        
        
        
# --- PROMPTS (UPDATED FOR PATH SAFETY) ---

TS_SYSTEM_PROMPT = """
You are a professional text summarization assistant specializing in accurately condensing written content 
while preserving key details and important information. Your task is to extract and present the most relevant points of a document,
including the title, author(s), affiliation(s),and other critical metadata, in a clear and concise manner. Your output must be formatted as a Markdown file.
"""

TS_USER_PROMPT = """
Here is a document for you to summarize: <Document text> 
Please summarize this document and generate a Markdown file. 
Ensure the summary includes the following:
1. The title of the document (if available).
2. The author's name(s) (if available).
3. The author's affiliation(s) (if mentioned).
4. A concise summary of the main content,focusing on key points, findings, or conclusions. 
Output the summary in a Markdown format as <TextSummary.md>. Ensure the details are accurate and well-organized.
"""

IC_SYSTEM_PROMPT = """
You are an expert image captioning assistant. 
Your role is to generate meaningful captions for images based on their content and the context provided in a document. 
Ensure that your captions are accurate, descriptive, and aligned with the references in the document text. 
Present your output in a clear and organized Markdown format.
"""

# UPDATED: Added instruction to reference specific filenames provided in the input
IC_USER_PROMPT = """
Here is the document text: <Document Text>
Here are the images (provided as visual input). 

Below is a list of the filenames corresponding to these images:
AVAILABLE IMAGES:
<File List>

Please analyze the images and, based on their content and the context of their references in the document:
1. Assign a title to each image.
2. Provide a detailed explanation of what the image shows and its relevance.
3. Indicate where the image is referenced in the text.
4. **CRITICAL:** When referring to the image, you MUST use the exact filename provided in the list above (e.g., "images/chart.png").
Output your captions in a Markdown file named <ImageCaption.md>.
"""

CG_SYSTEM_PROMPT = """
You are a highly skilled code generation assistant. Your role is to generate high-quality, well-structured code based on the provided instructions, 
ensuring it adheres to the specified requirements and formatting conventions. Your outputs should be accurate, organized, and easy to use.
"""

# UPDATED: Added Requirement #1 to force relative paths and forward slashes
# UPDATE THIS IN SlideGeneration.py
CG_USER_PROMPT = """
You are a JSON Architect. Your task is to structure slide data.
RULES:
1. The "content" list should ONLY contain text bullets.
2. The "image_path" should be a SEPARATE key outside the content list.
3. Use the exact filenames provided (e.g., "images/diagram1.png").

JSON SCHEMA:
[
  {
    "title": "Slide Title",
    "content": ["bullet 1", "bullet 2"],
    "image_path": "images/example.png"
  }
]
"""

CR_SYSTEM_PROMPT = """
You are a highly skilled code reviewer. Your role is to carefully analyze and evaluate code for correctness, clarity, 
and adherence to the given specifications. Your feedback should be precise, constructive, and well-structured.
"""

CR_USER_PROMPT = """
Here is the Slidev grammar specification: <SlidevGrammar> Additionally, here are two files: 
Text Summary: <TextSummary.md> 
Image Captions: <ImageCaption.md> 
Please review the code in <SlidevCode.md> to ensure:
1. It adheres to the Slidev grammar described in <Slidev Grammar>.
2. The content aligns with the information in <TextSummary.md> and <ImageCaption.md>.
3. **IMAGE PATH CHECK**: Verify that all image paths start with "images/" and do NOT use backslashes or absolute paths.
4. The code meets all content and layout requirements, including handling of images, text,and page structure as specified. 
Output your review as a Markdown file named <CodeReview.md>. Your review should clearly identify any errors or inconsistencies in the code, 
along with suggestions for improvement. If the code is correct, confirm that it meets all requirements.
"""

PR_SYSTEM_PROMPT = """
You are an expert visual page reviewer. Your role is to evaluate the layout and design of slides, 
ensuring they are visually appealing and properly aligned. Your feedback should be clear, actionable, 
and focused on improving the layout without altering the core content.
"""

PR_USER_PROMPT = """
Here are is a slide: <SlidevPages> 
Additionally, here are two supporting files: 
Image Information: <ImageCaption.md> 
Original Images: <Document Images> 
Please review each slide to check:
1. Whether any text or image exceeds the slide boundaries.
2. Whether the layout ensures a proper balance between text and images, avoiding over crowding or large empty spaces.
3. Whether the font sizes and styles are legible and consistent throughout the slide, ensuring readability without clashing with the visuals.
4. Whether the aspect ratios of images are preserved, and whether wide or tall images are placed appropriately without distorting the layout.
For each slide:
1. Indicate whether modifications are needed by answering with "yes" or "no".
2. If "yes", provide specific suggestions to adjust the positions of existing images. Do not add or remove any images. 
Output your review as a Markdown file named <PageReview.md>. Ensure your feedback is concise and easy to follow.
"""

# Also update the system prompt for the final agent
CGR_SYSTEM_PROMPT = "You are a specialized JSON data formatter. You never output Markdown or Slidev code. You only output structured JSON arrays."

# UPDATED: Enforce path rule in the correction phase too
# Force the final agent to stay in JSON mode
CGR_USER_PROMPT = """
Act as a JSON Architect. Convert the following text and image summaries into a valid JSON array for PowerPoint.

INPUTS:
Text Summary: <TextSummary.md> 
Image Captions: <ImageCaption.md> 
Review Feedback: <CodeReview.md> or <PageReview.md> 

CRITICAL RULES:
1. OUTPUT ONLY A VALID JSON ARRAY.
2. DO NOT include markdown backticks (```json).
3. DO NOT include headers like "# SlidevCode.md".
4. START with '[' and END with ']'.

SCHEMA:
[
  {
    "title": "Slide Title",
    "layout": "bullet_points",
    "content": ["point 1", "point 2"],
    "image_path": "images/filename.png"
  }
]
"""

# --- SLIDEV TEMPLATES ---
SLIDEV_GRAMMAR = """
# Slidev Markdown Grammar
---
# Frontmatter (first slide)
layout: cover
title: My Slide Title
---

# Standard Slide

# Slide Title

- Bullet point 1
- Bullet point 2

---
layout: two-cols

::left::
# Left Column
Content here

::right::
# Right Column
Content here

---
# Image Slide
![Alt Text](images/filename.png)
"""

SLIDEV_PAGES = """
Available Layouts:
1. 'cover': Use for the first slide. Contains title and subtitle.
2. 'default': Standard Title + Content.
3. 'two-cols': Split screen. Use ::left:: and ::right:: separators.
4. 'image-right': Content on left, image on right.
5. 'center': Centered text for quotes or impact statements.
"""