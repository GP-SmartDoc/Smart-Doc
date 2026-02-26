from pptx import Presentation
import json
import os
import re
import copy

# [cite_start]Mapping to match the order of your layouts.pptx [cite: 12]
LAYOUT_MAP = {
    "title_subtitle": 0, "title_content": 1, "three_column": 2,
    "content_image": 3, "large_image": 4, "title_only": 5
}

# PowerPoint XML Namespaces
NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
}

def format_value(val):
    """Extracts values from dicts (no keys) and merges lists into bullets."""
    if isinstance(val, dict):
        merged = []
        for v in val.values():
            if isinstance(v, list):
                merged.extend([str(item) for item in v])
            else:
                merged.append(str(v))
        return "\n".join([f"• {p}" if not p.startswith('•') else p for p in merged])
    if isinstance(val, list):
        return "\n".join([f"• {item}" for item in val])
    return str(val) if val else ""

def scrub_xml_for_repair(element):
    """
    Recursively cleans XML of Relationship IDs AND Extension Lists.
    This is the ONLY way to keep animations without repair errors.
    """
    for node in element.iter():
        # 1. Remove Relationship attributes (r:id, r:embed)
        for attr in list(node.attrib):
            if "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}" in attr:
                del node.attrib[attr]
        
        # 2. Remove Extension Lists (extLst) - these store unique slide IDs 
        # that cause the Repair prompt in animations/transitions.
        if node.tag.endswith('extLst'):
            parent = node.getparent()
            if parent is not None:
                parent.remove(node)
    return element

def save_as_pptx(raw_llm_output, template_path="layouts.pptx", output_path="generated_slides.pptx"):
    if not os.path.exists(template_path): return
    
    prs = Presentation(template_path)
    template_count = len(prs.slides)
    
    match = re.search(r'\[\s*\{.*\}\s*\]', raw_llm_output, re.DOTALL)
    if not match: return
    slides_data = json.loads(match.group(0))

    for data in slides_data:
        layout_key = data.get("layout_name", data.get("layout", "title_content"))
        idx = LAYOUT_MAP.get(layout_key, 1)
        
        # 1. Add fresh slide and CLEAR default placeholders
        slide_layout = prs.slide_masters[0].slide_layouts[idx]
        new_slide = prs.slides.add_slide(slide_layout)
        for shp in list(new_slide.shapes):
            new_slide.shapes._spTree.remove(shp._element)
        
        # 2. Get template slide instance (where your animations live)
        template_slide = prs.slides[idx]
        
        # 3. Copy Shapes with Scrubbing
        for shape in template_slide.shapes:
            clean_shape_el = scrub_xml_for_repair(copy.deepcopy(shape._element))
            new_slide.shapes._spTree.append(clean_shape_el)

        # 4. CRITICAL: Copy Transitions and Animations with DEEP SCRUB
        template_el = template_slide._element
        for tag in ['transition', 'timing']:
            node = template_el.find(f'./p:{tag}', namespaces=NS)
            if node is not None:
                # Scrubbing the timing/transition node prevents the Repair error
                # by removing stale slide-ID references inside the XML.
                clean_node = scrub_xml_for_repair(copy.deepcopy(node))
                new_slide._element.append(clean_node)

        # 5. Tag Replacement (Case-Insensitive)
        content_source = {str(k).lower(): v for k, v in data.items() if k not in ["tags", "layout", "layout_name"]}
        if "tags" in data and isinstance(data["tags"], dict):
            content_source.update({str(k).lower(): v for k, v in data["tags"].items()})

        for shape in new_slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key in ["title", "subtitle", "content", "col1", "col2", "col3"]:
                            pattern = re.compile(r'\{\{' + re.escape(key) + r'\}\}', re.IGNORECASE)
                            if pattern.search(run.text):
                                run.text = pattern.sub(format_value(content_source.get(key)), run.text)
                        if "{{" in run.text:
                            run.text = re.sub(r'\{\{.*?\}\}', '', run.text)

        # 6. Image Replacement (Targets DUMMY placeholders)
        img_path = data.get("image_path", data.get("image"))
        if img_path and os.path.exists(img_path):
            for shape in list(new_slide.shapes):
                sh_name = (shape.name or "").upper()
                if "DUMMY" in sh_name or (shape.has_text_frame and "DUMMY" in shape.text.upper()):
                    l, t, w, h = shape.left, shape.top, shape.width, shape.height
                    new_slide.shapes._spTree.remove(shape._element)
                    new_slide.shapes.add_picture(img_path, l, t, w, h)
                    break

    # 7. Cleanup template slides from the output file
    for _ in range(template_count):
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    prs.save(output_path)